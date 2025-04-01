import os
import requests
import pandas as pd
from openai import OpenAI
import json
import logging
import traceback
from datetime import datetime
from dotenv import load_dotenv
import streamlit as st
from io import BytesIO
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from requests.adapters import HTTPAdapter
from urllib3.util import Retry
import concurrent.futures
import graphviz
import tempfile
import uuid
from typing import List, Dict, Optional, Any, Tuple, Union
from dataclasses import dataclass
from enum import Enum
import numpy as np
import openai
import re
import textwrap
from api_key_manager import APIKeyManager
from exceptions import ConfigurationError, PatentAnalysisError

# Load environment variables first
load_dotenv()

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

# 상수 정의
class Constants:
    GPT_MODEL = "gpt-3.5-turbo"
    GPT_MODEL_4 = "gpt-4"
    DEFAULT_ERROR_MESSAGE = "분석 실패"
    DEFAULT_IMAGE_SIZE = "1024x1024"
    DEFAULT_IMAGE_QUALITY = "standard"
    DEFAULT_NUM_IMAGES = 1
    DEFAULT_SIMILARITY_THRESHOLD = 3.0
    DEFAULT_NUM_RESULTS = 20
    DEFAULT_TOP_N = 5

class PatentType(Enum):
    HARDWARE = "hardware_component"
    SYSTEM = "integrated_system"
    ALGORITHM = "control_algorithm"
    INTERFACE = "user_interface"
    PROCESS = "manufacturing_process"
    MATERIAL = "new_material"
    SENSOR = "sensor_technology"

@dataclass
class PatentDomain:
    primary_domain: str
    secondary_domains: List[str]
    idea_type: PatentType
    technical_approach: str

class ConfigurationError(Exception):
    """설정 관련 오류를 위한 사용자 정의 예외"""
    pass

class PatentAnalysisError(Exception):
    """특허 분석 관련 오류를 위한 사용자 정의 예외"""
    pass

class ErrorHandler:
    """공통 오류 처리를 위한 클래스"""
    
    @staticmethod
    def handle_api_error(func):
        """API 호출 관련 오류 처리를 위한 데코레이터"""
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except requests.exceptions.RequestException as e:
                logging.error(f"API 요청 오류: {str(e)}")
                raise PatentAnalysisError(f"API 요청 중 오류가 발생했습니다: {str(e)}")
            except Exception as e:
                logging.error(f"예상치 못한 오류: {str(e)}")
                traceback.print_exc()
                raise PatentAnalysisError(f"처리 중 오류가 발생했습니다: {str(e)}")
        return wrapper
    
    @staticmethod
    def handle_file_operation(func):
        """파일 작업 관련 오류 처리를 위한 데코레이터"""
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except IOError as e:
                logging.error(f"파일 작업 오류: {str(e)}")
                raise PatentAnalysisError(f"파일 작업 중 오류가 발생했습니다: {str(e)}")
            except Exception as e:
                logging.error(f"예상치 못한 오류: {str(e)}")
                traceback.print_exc()
                raise PatentAnalysisError(f"처리 중 오류가 발생했습니다: {str(e)}")
        return wrapper

class CacheManager:
    """캐시 관리를 위한 클래스"""
    
    def __init__(self, cache_dir: str = "cache"):
        self.cache_dir = cache_dir
        os.makedirs(cache_dir, exist_ok=True)
    
    @ErrorHandler.handle_file_operation
    def get_cached_data(self, key: str, max_age_hours: int = 24) -> Optional[Dict]:
        """캐시된 데이터 조회"""
        cache_file = os.path.join(self.cache_dir, f"{key}.json")
        
        if not os.path.exists(cache_file):
            return None
            
        # 캐시 유효기간 확인
        file_age = datetime.now().timestamp() - os.path.getmtime(cache_file)
        if file_age > (max_age_hours * 3600):
            return None
            
        with open(cache_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    @ErrorHandler.handle_file_operation
    def save_to_cache(self, key: str, data: Dict) -> None:
        """데이터를 캐시에 저장"""
        cache_file = os.path.join(self.cache_dir, f"{key}.json")
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

# HTTP 세션 설정
session = requests.Session()
retry_strategy = Retry(
    total=5,
    backoff_factor=2,
    status_forcelist=[429, 500, 502, 503, 504],
    allowed_methods=["HEAD", "GET", "POST", "OPTIONS"]
)
adapter = HTTPAdapter(max_retries=retry_strategy, pool_connections=10, pool_maxsize=10)
session.mount("http://", adapter)
session.mount("https://", adapter)
session.timeout = (10, 60)

class PatentAnalysisError(Exception):
    """Custom exception for patent analysis errors"""
    pass

class EnhancedPatentAnalysisAgent:
    """Enhanced patent analysis agent with improved similarity calculation"""
    
    def __init__(self):
        """Initialize the agent with API key"""
        self.api_key = APIKeyManager.get_openai_api_key()
        self.client = OpenAI(api_key=self.api_key)
        self.cache_manager = CacheManager()
        self.domain_info = None
        self.semantic_keywords = None
        self.idea_description = None
        self.ge_categories = self.get_ge_categories()

    def get_ge_categories(self):
        """GE Appliances 카테고리 구조 로드"""
        return {
            "PEB_Categories": [
                "1. Refrigeration",
                "2. Cooking Products",
                "3. Dishwashers",
                "4. Home Laundry",
                "5. Air & Water Products",
                "6. Power Electronics & Connected Home",
                "7. Small Appliances/FirstBuild",
                "8. Global Specialty Products"
            ],
            "Full_Classification_Tree": {
                "1. Refrigeration": {
                    "1.1 - Refrigerators": [
                        "1.1.1 - Compressors/Refrigerants/Circuits",
                        "1.1.2 - Heat Exchangers",
                        "1.1.3 - Fans & Air Flow Systems",
                        "1.1.4 - Defrost Systems",
                        "1.1.5 - Control Algorithms",
                        "1.1.6 - F&A",
                        "1.1.7 - Doors & Structures",
                        "1.1.8 - Ice & Water (Ice makers, system config., ice and water level sensing, etc..)",
                        "1.1.9 - Non Vapor Compression Cooling"
                    ],
                    "1.2 - Freezers": [
                        "1.2.1 - Chest Freezers",
                        "1.2.2 - Upright Freezers"
                    ],
                    "1.3 - Water Filtration": [
                        "1.3.1 - Interfaces",
                        "1.3.2 - Filtration Media",
                        "1.3.3 - Water Treatment (Softening, Ionization, etc..)",
                        "1.3.4 - Quality Detection & Connectivity (leak, lead, heavy metals, impurities)",
                        "1.3.5 - Dispensers & Softeners"
                    ],
                    "1.4 - Ice Makers": [
                        "1.4.1 - Countertop",
                        "1.4.2 - Under Counter",
                        "1.4.3 - Built-In",
                        "1.4.4 - Other"
                    ],
                    "1.5 - Electronics": [
                        "1.5.1 - Electronics",
                        "1.5.2 - Interoperability",
                        "1.5.3 - Machine Learning"
                    ]
                },
                "2. Cooking Products": {
                    "2.1 - Cooking Ranges": [
                        "2.1.1 - Gas Cooking Ranges",
                        "2.1.2 - Electric Cooking Ranges",
                        "2.1.3 - Dual Fuel Cooking Ranges",
                        "2.1.4 - Induction Cooking Ranges",
                        "2.1.5 - Cooking Ranges in General"
                    ],
                    "2.2 - Ovens (Wall Ovens)": [
                        "2.2.1 - Gas Ovens",
                        "2.2.2 - Electric Ovens",
                        "2.2.3 - Ovens in General"
                    ],
                    "2.3 - Cooktops": [
                        "2.3.1 - Gas Cooktops",
                        "2.3.2 - Electric Cooktops",
                        "2.3.3 - Induction Cooktops",
                        "2.3.4 - Cooktops in General"
                    ],
                    "2.4 - Warming Drawers": [],
                    "2.5 - Electronics": [
                        "2.5.1 - Electronics",
                        "2.5.2 - Interoperability",
                        "2.5.3 - Machine Learning"
                    ]
                },
                "3. Dishwashers": {
                    "3.1 - Dishwashers": [
                        "3.1.1 - Wash, Dry, Racking",
                        "3.1.2 - Tub Structure, Door, Sound",
                        "3.1.3 - Electronics and Controls"
                    ],
                    "3.2 - Disposers & Compactors": [
                        "3.2.1 - Disposers",
                        "3.2.2 - Compactors"
                    ],
                    "3.3 - Electronics": [
                        "3.3.1 - Electronics",
                        "3.3.2 - Interoperability",
                        "3.3.3 - Machine Learning"
                    ]
                },
                "4. Home Laundry": {
                    "4.1 - Top Load Washers": [
                        "4.1.1 - Wash Systems",
                        "4.1.2 - Drive System",
                        "4.1.3 - Suspension/Vibration",
                        "4.1.4 - Fill/Dispense/Drain",
                        "4.1.5 - Tub/Drum",
                        "4.1.6 - Structure/Control Panel",
                        "4.1.7 - Electronics/Sensing/Smart Home"
                    ],
                    "4.2 - Front Load Washers": [
                        "4.2.1 - Wash Systems",
                        "4.2.2 - Drive System",
                        "4.2.3 - Suspension/Vibration",
                        "4.2.4 - Fill/Dispense/Drain",
                        "4.2.5 - Tub/Drum",
                        "4.2.6 - Structure/Control Panel",
                        "4.2.7 - Electronics/Sensing/Smart Home"
                    ],
                    "4.3 - Dryers": [
                        "4.3.1 - Heating System",
                        "4.3.2 - Airflow System",
                        "4.3.3 - Lint Management",
                        "4.3.4 - Sealed System",
                        "4.3.5 - Structure/Control Panel",
                        "4.3.6 - Electronics/Sensing/Smart Home" 
                    ],
                    "4.4 - Combination Washer/Dryers": [
                        "4.4.1 - Wash Systems",
                        "4.4.2 - Drive System",
                        "4.4.3 - Suspension/Vibration",
                        "4.4.4 - Fill/Dispense/Drain",
                        "4.4.5 - Tub/Drum",
                        "4.4.6 - Heating System",
                        "4.4.7 - Airflow System",
                        "4.4.8 - Lint Management",
                        "4.4.9 - Sealed System",
                        "4.4.10 - Structure/Control Panel",
                        "4.4.11 - Electronics/Sensing/Smart Home"
                    ],
                    "4.5 - Accessories": []
                },
                "5. Air & Water Products": {
                    "5.1 - Air Conditioners": [
                        "5.1.1 - Window Air Conditioners",
                        "5.1.2 - Split Air Conditioners",
                        "5.1.3 - Air Conditioners/Heat Pumps in General",
                        "5.1.4 - Commercial Air Conditioners (PTACs, Vertical, Built-in)",
                        "5.1.5 - Large Commercial Central Air Conditioners",
                        "5.1.6 - Recreational Vehicle Air Conditioners",
                        "5.1.7 - Room Thermostats",
                        "5.1.8 - Air Conditioning Accessories (Sleeves, Assembly Aids, Ducting, etc.)",
                        "5.1.9 - Air Conditioning Remote Connected Strategies",
                        "5.1.10 - Electronics",
                        "5.1.11 - Interoperability",
                        "5.1.12 - Machine Learning"
                    ],
                    "5.2 - Dehumidifiers & Air Purifiers": [
                        "5.2.1 - Dehumidifiers",
                        "5.2.2 - Air Purifiers",
                        "5.2.3 - Electronics",
                        "5.2.4 - Interoperability",
                        "5.2.5 - Machine Learning"
                    ],
                    "5.3 - Water Heaters": [
                        "5.3.1 - Gas Water Heaters",
                        "5.3.2 - Electric Water Heaters",
                        "5.3.3 - Heat Pump Water Heaters",
                        "5.3.4 - Electronics",
                        "5.3.5 - Interoperability",
                        "5.3.6 - Machine Learning"
                    ]
                },
                "6. Power Electronics & Connected Home": {
                    "6.1 - Home Energy Manager": [
                        "6.1.1 - Communication",
                        "6.1.2 - HEM Module",
                        "6.1.3 - Power Management"
                    ],
                    "6.2 - Home Networking": [
                        "6.2.1 - Internet of Things",
                        "6.2.2 - Interoperability between Appliances",
                        "6.2.3 - Remote Control, Data Acquisition, Diagnostics and Analytics (AI/Machine Learning on Data to Improve)"
                    ],
                    "6.3 - Electronics (Covering Multiple Appliances)": [
                        "6.3.1 - Embedded Software",
                        "6.3.2 - Power Management/Power Conversion (Inverters and Power Supply)",
                        "6.3.3 - Remote Server",
                        "6.3.4 - Induction (Not Related to Cooking Technologies)",
                        "6.3.5 - Controls (Generic Classification of Appliances Controls)",
                        "6.3.6 - Motor Control (Operation/Control of Variable Speed Drives)"
                    ]
                },
                "7. Small Appliances/FirstBuild": {
                    "7.1 - FirstBuild": [],
                    "7.2 - General Application": []
                },
                "8. Global Specialty Products": {},
                "8. Connected Home": {},
                "9. Microwave Ventilation & Hoods": {
                    "9.1 - Microwave": [],
                    "9.2 - Ventilation": [],
                    "9.3 - Hoods": []
                }
            }
        }
    

    @ErrorHandler.handle_api_error
    def search_patents(self, idea_description, num_results=Constants.DEFAULT_NUM_RESULTS):
        """Perform patent search"""
        print(f"Searching patents related to: {idea_description}")
        
        # Store as class attribute (used for relevance score calculation)
        self.idea_description = idea_description
        
        # Check cache
        cache_key = f"patents_{hash(idea_description)}"
        cached_data = self.cache_manager.get_cached_data(cache_key)
        if cached_data:
            print("Loading cached results...")
            return cached_data
        
        try:
            # Extract keywords
            self.semantic_keywords = self._extract_semantic_keywords(idea_description)
            self.domain_info = self._classify_patent_domain(idea_description)
            
            # Generate search terms
            search_terms = []
            if isinstance(self.semantic_keywords, dict) and "all" in self.semantic_keywords:
                search_terms.extend(self.semantic_keywords["all"][:5])
            elif isinstance(self.semantic_keywords, list):
                search_terms.extend(self.semantic_keywords[:5])
                
            if self.domain_info and "primary_domain" in self.domain_info:
                search_terms.insert(0, self.domain_info["primary_domain"])
                
            # Optimize search query
            search_query = " ".join(search_terms[:3])
            print(f"Search query: {search_query}")
            
            # Use OpenAI directly for search
            prompt = f"""
            Search for patents related to: {search_query}
            
            Return patent information in a structured format with these fields:
            - title: Patent Title
            - patent_id: Patent ID/Number
            - publication_date: Publication Date
            - inventors: List of inventors
            - assignee: Assignee Name
            - abstract: Patent Abstract
            - cpc_classifications: CPC Classifications
            - link: Patent URL
            
            Limit to {num_results} patents maximum.
            Format your response as a JSON array.
            """
            
            response = self.client.chat.completions.create(
                model="gpt-4o-search-preview",
                messages=[
                    {
                        "role": "system", 
                        "content": "You are a patent search expert. Return patent information in a structured JSON format."
                    },
                    {"role": "user", "content": prompt}
                ],
                web_search_options={
                    "search_context_size": "high"
                },
                max_tokens=4000
            )
            
            # Parse response
            content = response.choices[0].message.content
            content = content.strip()
            
            # Find JSON array
            match = re.search(r'\[.*\]', content, re.DOTALL)
            if match:
                json_str = match.group(0)
                patents = json.loads(json_str)
            else:
                match = re.search(r'\{.*\}', content, re.DOTALL)
                if match:
                    json_str = match.group(0)
                    patents_obj = json.loads(json_str)
                    patents = patents_obj.get("patents", [patents_obj])
                else:
                    patents = []
            
            if not patents:
                print("No search results found. Generating mock data...")
                return self.generate_mock_patent_data(idea_description, num_results)
            
            # Cache results
            self.cache_manager.save_to_cache(cache_key, patents)
            return patents
            
        except Exception as e:
            print(f"Patent search failed: {e}")
            traceback.print_exc()
            return self.generate_mock_patent_data(idea_description, num_results)
    
    def generate_mock_patent_data(self, idea_description, num_results=5):
        """모의 특허 데이터 생성 - 개발 및 테스트용"""
        logging.warning("실제 특허 검색 실패로 인한 모의 데이터 생성")
        
        try:
            domain_info = self._classify_patent_domain(idea_description)
            primary_domain = domain_info.get("primary_domain", "electronics")
            
            mock_patents = []
            base_relevance = 80.0  # 기본 관련성 점수
            
            for i in range(num_results):
                # 관련성 점수 감소
                relevance_score = max(30.0, base_relevance - (i * 10))
                
                # 특허 제목 생성
                if i % 3 == 0:
                    title = f"지능형 {primary_domain} 시스템 및 제어 방법"
                elif i % 3 == 1:
                    title = f"사용자 인식 기반 {primary_domain} 장치"
                else:
                    title = f"AI 기반 {primary_domain} 최적화 시스템"
                
                # 모의 특허 데이터 생성
                mock_patents.append({
                    'title': title,
                    'patent_id': f"MOCK-{i+1}",
                    'source': '테스트 데이터 (서비스 일시 중단)',
                    'url': "#",
                    'description': (
                        f"이 특허는 {primary_domain} 분야에서 발생하는 문제를 해결하기 위한 "
                        f"혁신적인 접근 방식을 제시합니다. 제안된 시스템은 사용자의 요구사항을 "
                        f"효율적으로 처리하며, 에너지 효율성과 사용자 편의성을 극대화합니다."
                    ),
                    'inventors': "홍길동, 김철수",
                    'assignee': f"{primary_domain.title()} 기술 주식회사",
                    'relevance_score': relevance_score,
                    'domain': primary_domain,
                    'publication_date': "2024-01-01",
                    'relevant_ge_categories': self._determine_relevant_categories(idea_description)
                })
            
            return mock_patents
            
        except Exception as e:
            logging.error(f"모의 데이터 생성 중 오류 발생: {str(e)}")
            # 최소한의 데이터 반환
            return [{
                'title': f"테스트 특허 {i+1}",
                'patent_id': f"TEST-{i+1}",
                'source': '오류 복구 데이터',
                'relevance_score': 50.0,
                'domain': 'technology'
            } for i in range(num_results)]
    
    def _calculate_relevance_based_on_keywords(self, text, keywords):
        """특허 텍스트와 키워드 간의 관련성 점수 계산"""
        if not text or not keywords:
            return 0.0
        
        text = text.lower()
        score = 0.0
        
        for keyword in keywords:
            keyword = keyword.lower()
            if keyword in text:
                score += 1.0
                # 정확한 단어 매칭에 대해 추가 점수
                if re.search(r'\b' + re.escape(keyword) + r'\b', text):
                    score += 0.5
        
        # 정규화: 키워드 수로 나누고 최대 점수를 5.0으로 설정
        normalized_score = min(5.0, score / max(1, len(keywords)) * 5.0)
        return normalized_score
    
    def _determine_relevant_categories(self, idea_description):
        """Identify relevant GE Appliances categories for the patent idea"""
        # Format categories for AI prompt
        formatted_categories = json.dumps(self.ge_categories, indent=2)
        
        prompt = f"""
        Based on the following patent idea description, identify the most relevant GE Appliances categories.
        Return the 1-3 most relevant categories from both the PEB_Categories list and the Full_Classification_Tree.
        
        Patent idea description:
        {idea_description}
        
        GE Appliances Categories:
        {formatted_categories}
        
        Return a JSON object with two properties:
        1. "peb_category": The single most relevant PEB category
        2. "full_categories": Array of 1-3 most relevant full library classifications (use the full path like "1. Refrigeration > 1.1 - Refrigerators > 1.1.5 - Control algorithms")
        """
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
            
            result = json.loads(response.choices[0].message.content)
            return result
        except Exception as e:
            print(f"Error determining categories: {e}")
            return {
                "peb_category": "Refrigeration",  # Default
                "full_categories": []
            }
    
    def _extract_semantic_keywords(self, idea_description):
        prompt = f"""
        Extract semantic keywords from the following patent idea description.
        Focus on understanding the core technology, problem being solved, and unique approach.
        Extract 3 types of keywords:
    
        1. Core Technical Concepts (3-4 keywords): The fundamental technical principles or mechanisms
        2. Problem Domain Keywords (2-3 keywords): What problem is being solved
        3. Unique Approach Keywords (2-3 keywords): What makes this solution unique
        4. Alternative Terminology (2-3 keywords): Different ways people might refer to similar concepts
    
        Return the keywords as a JSON object with these four categories as properties, each with an array of strings.
    
        Patent idea description:
        {idea_description}
        """
    
        try:
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
        
            result = json.loads(response.choices[0].message.content)
        
        # 결과 유효성 검사 및 기본값 설정
            if not isinstance(result, dict):
                result = {}
            
        # 모든 키워드를 단일 리스트로 통합 (유효성 검사 추가)
            all_keywords = []
            for category, keywords in result.items():
                if isinstance(keywords, list):
                    all_keywords.extend(keywords)
            
            return {
                "structured": result,
                "all": all_keywords
            }
        except Exception as e:
            print(f"Error extracting semantic keywords: {e}")
            return {
                "structured": {},
                "all": []
            }
    
    def _classify_patent_domain(self, idea_description):
        """Classify the patent domain and type for better targeting"""
        prompt = f"""
        Classify the following patent idea into its domain and type.
        
        Patent idea description:
        {idea_description}
        
        Return a JSON object with these properties:
        1. "primary_domain": One of [refrigeration, cooking, dishwashing, laundry, air_conditioning, water_treatment, electronics, smart_home, small_appliances]
        2. "secondary_domains": Array of 0-2 additional relevant domains
        3. "idea_type": One of [hardware_component, integrated_system, control_algorithm, user_interface, manufacturing_process, new_material, sensor_technology]
        4. "technical_approach": Brief phrase describing the technical approach (e.g., "machine learning for power optimization")
        """
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
            
            result = json.loads(response.choices[0].message.content)
            return result
        except Exception as e:
            print(f"Error classifying patent domain: {e}")
            return {
                "primary_domain": "electronics",
                "secondary_domains": [],
                "idea_type": "integrated_system",
                "technical_approach": "automated control"
            }
    
    def _create_semantic_search_queries(self, idea_description, semantic_keywords, domain, idea_type):
        """Create diverse search queries with semantic understanding"""
        # Extract keywords from the structured semantic keywords
        core_keywords = semantic_keywords["structured"].get("Core Technical Concepts", [])[:3]
        problem_keywords = semantic_keywords["structured"].get("Problem Domain Keywords", [])[:2]
        approach_keywords = semantic_keywords["structured"].get("Unique Approach Keywords", [])[:2]
        alt_keywords = semantic_keywords["structured"].get("Alternative Terminology", [])[:2]
        
        # Create diverse queries
        queries = []
        
        # Main query with core concepts and domain
        if core_keywords:
            main_query = f'"{core_keywords[0]}" "{core_keywords[1]}" {domain}'
            queries.append(main_query)
        
        # Problem-focused query with technical terms
        if problem_keywords and core_keywords:
            problem_query = f'"{problem_keywords[0]}" "{core_keywords[0]}" patent'
            queries.append(problem_query)
        
        # Approach-focused query with technical details
        if approach_keywords and core_keywords:
            approach_query = f'"{approach_keywords[0]}" "{core_keywords[0]}" {idea_type}'
            queries.append(approach_query)
        
        # Alternative terminology query with domain context
        if alt_keywords:
            alt_query = f'"{alt_keywords[0]}" {domain} {idea_type} patent'
            queries.append(alt_query)
        
        # Add fallback query if needed
        if len(queries) < 3:
            # Create a more focused query from the idea description
            words = idea_description.split()
            if len(words) >= 4:
                fallback_query = f'"{words[0]} {words[1]}" "{words[2]} {words[3]}" patent'
                queries.append(fallback_query)
        
        # Add domain-specific queries
        if domain:
            domain_query = f'"{domain}" "{core_keywords[0] if core_keywords else ""}" patent'
            queries.append(domain_query)
        
        return queries
    
    def _calculate_semantic_relevance(self, title, description, semantic_keywords, idea_description, domain_info):
        """Calculate semantic relevance between patent and idea"""
        score = 0.0
        
        # Core keyword matching with weighting
        if isinstance(semantic_keywords, dict) and "all" in semantic_keywords:
            keywords = semantic_keywords["all"]
        else:
            keywords = semantic_keywords if isinstance(semantic_keywords, list) else []
        
        for keyword in keywords:
            if keyword.lower() in title.lower():
                score += 3.0
            if keyword.lower() in description.lower():
                score += 1.0
        
        # Domain relevance scoring
        primary_domain = domain_info.get("primary_domain", "")
        if primary_domain and primary_domain.lower() in (title.lower() + " " + description.lower()):
            score += 2.0
            
        for secondary_domain in domain_info.get("secondary_domains", []):
            if secondary_domain.lower() in (title.lower() + " " + description.lower()):
                score += 1.0
        
        # Idea type matching
        idea_type = domain_info.get("idea_type", "")
        if idea_type and idea_type.lower() in (title.lower() + " " + description.lower()):
            score += 1.5
            
        # Technical approach matching
        technical_approach = domain_info.get("technical_approach", "")
        if technical_approach and technical_approach.lower() in description.lower():
            score += 2.5
        
        # 정규화: 최대 점수를 100으로 설정
        normalized_score = min(100.0, score * 10)
        return normalized_score
    
    @ErrorHandler.handle_api_error
    def get_patent_details(self, patent_id):
        if not patent_id:
            return {}
        
        try:
            # Check cache
            cache_key = f"patent_details_{patent_id}"
            cached_data = self.cache_manager.get_cached_data(cache_key)
            if cached_data:
                return cached_data
        
            # Use OpenAI to get patent information instead of API request
            prompt = f"Find detailed information for patent number {patent_id}. Include title, abstract, inventors, assignee, filing date, claims, and other relevant information."
        
            response = self.client.chat.completions.create(
                model="gpt-4o-search-preview",
                messages=[{"role": "user", "content": prompt}],
                web_search_options={"search_context_size": "high"}
            )
        
            # Parse response
            patent_data = self._parse_patent_details(response.choices[0].message.content)
        
            # Cache results
            self.cache_manager.save_to_cache(cache_key, patent_data)
        
            return patent_data
    
        except Exception as e:
            logging.error(f"Failed to retrieve patent details: {e}")
            return {}
    
    @ErrorHandler.handle_api_error
    def get_patents_details(self, patents, max_workers=5):
        """Retrieve detailed information for multiple patents in parallel"""
        print(f"Retrieving details for {len(patents)} patents in parallel...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit tasks for each patent
            future_to_patent = {
                executor.submit(self.get_patent_details, patent['patent_id']): 
                (i, patent) for i, patent in enumerate(patents)
            }
            
            # Process completed tasks
            completed = 0
            for future in concurrent.futures.as_completed(future_to_patent):
                i, patent = future_to_patent[future]
                try:
                    details = future.result()
                    if details:
                        # Merge details with original patent information
                        patents[i].update(details)
                        
                        # Calculate relevance score
                        relevance_score = self._calculate_relevance(
                            patents[i].get('title', '') + ' ' + patents[i].get('description', ''),
                            keywords=getattr(self, 'semantic_keywords', None),
                            domain_info=getattr(self, 'domain_info', None)
                        )
                        patents[i]['relevance_score'] = relevance_score
                    
                    completed += 1
                    print(f"Progress: {completed}/{len(patents)} patents processed")
                    
                except Exception as e:
                    logging.error(f"Failed to retrieve details for patent {patent['patent_id']}: {e}")
                    traceback.print_exc()
        
        return patents
    
    def analyze_patents(self, idea_description, patents_list, top_n=5):
        """Perform patent analysis"""
        if not patents_list:
            return []
            
        try:
            # Classify domain and determine categories
            domain_info = self._classify_patent_domain(idea_description)
            categories = self._determine_relevant_categories(idea_description)
            
            # Analyze each patent
            analyzed_patents = []
            for patent in patents_list:
                analysis = self._deep_semantic_comparison(
                    idea_description,
                    patent.get('title', ''),
                    patent.get('description', ''),
                    patent.get('claims', [])
                )
                
                # Add analysis results to patent data
                patent['analysis'] = analysis
                patent['similarity_score'] = analysis.get('similarity_score', 0)
                analyzed_patents.append(patent)
            
            # Sort by similarity score
            analyzed_patents.sort(key=lambda x: x['similarity_score'], reverse=True)
            
            # Return top N results
            return analyzed_patents[:top_n]
            
        except Exception as e:
            print(f"Error during patent analysis: {e}")
            return []
    
    def _deep_semantic_comparison(self, idea_description, patent_title, patent_content, patent_claims=None):
        """Perform deep semantic comparison between user idea and patent"""
        try:
            prompt = f"""
            Perform a deep semantic comparison between the user's patent idea and an existing patent.
            
            [User Idea]
            {idea_description}
            
            [Existing Patent]
            Title: {patent_title}
            Content: {patent_content[:2000]}  # Limit content length
            
            {f"Claims: {patent_claims[:3]}" if patent_claims else ""}
            
            Analyze:
            1. Overall Similarity (0-100 scale): How similar the core technology and approach are
            2. Key Similarities (exact 3): Specific technical similarities in approach, mechanisms, or solutions
            3. Key Differences (exact 3): Specific technical differences that distinguish the user idea
            4. Key Features (exact 3): The most important features of the existing patent
            
            Return analysis as a JSON object with these four properties.
            Make sure similarities and differences focus on technical aspects, not generalities.
            """
            
            response = self.client.chat.completions.create(
                model=Constants.GPT_MODEL,
                messages=[{"role": "user", "content": prompt}]
            )
            
            analysis = json.loads(response.choices[0].message.content)
            return analysis
            
        except Exception as e:
            print(f"Error during semantic comparison: {e}")
            return {
                "similarity_score": 0,
                "similarities": [Constants.DEFAULT_ERROR_MESSAGE],
                "differences": [Constants.DEFAULT_ERROR_MESSAGE],
                "key_features": [Constants.DEFAULT_ERROR_MESSAGE]
            }
    
    def generate_comparison_visual(self, idea_description, analyzed_patents):
        """Generate a visual comparison chart of patent similarities"""
        if not analyzed_patents:
            return None
        
        try:
            # Create data for visualization
            titles = [p["title"][:30] + "..." if len(p["title"]) > 30 else p["title"] for p in analyzed_patents]
            scores = [p.get("similarity_score", 0) for p in analyzed_patents]
            
            # Generate the chart
            fig, ax = plt.subplots(figsize=(10, 6))
            bars = ax.barh(titles, scores, color='royalblue')
            
            # Add similarity percentage at the end of each bar
            for i, bar in enumerate(bars):
                width = bar.get_width()
                ax.text(width + 1, bar.get_y() + bar.get_height()/2, 
                       f'{width:.1f}%', ha='left', va='center')
            
            ax.set_xlabel('Similarity Score (%)')
            ax.set_title('Patent Similarity Comparison')
            ax.set_xlim(0, 100)
            
            # Save to a bytes buffer
            buf = io.BytesIO()
            fig.savefig(buf, format='png', bbox_inches='tight')
            buf.seek(0)
            
            return buf
            
        except Exception as e:
            print(f"Error generating comparison visual: {e}")
            return None
    
    def recommend_differentiation(self, idea_description, analyzed_patents):
        """Generate differentiation strategies with concrete technical suggestions"""
        if not analyzed_patents:
            return {}
            
        try:
            # Extract key information from analyzed patents
            patent_summaries = []
            for patent in analyzed_patents:
                analysis = patent.get('analysis', {})
                patent_summaries.append({
                    "title": patent.get("title", ""),
                    "similarities": analysis.get("similarities", []),
                    "differences": analysis.get("differences", []),
                    "key_features": analysis.get("key_features", []),
                    "similarity_score": patent.get("similarity_score", 0)
                })
            
            prompt = f"""
            Based on the user's patent idea and analysis of similar patents, provide concrete technical differentiation strategies.
            
            [User Idea]
            {idea_description}
            
            [Similar Patent Analysis]
            {json.dumps(patent_summaries, ensure_ascii=False)}
            
            Provide:
            1. Key competing patents (top 3 most similar)
            2. Technical differentiation opportunities (5 specific technical approaches that would make this unique)
            3. Improvement suggestions (5 specific technical enhancements to strengthen the patent)
            4. Patentability strategies (3 approaches to strengthen patent claims)
            5. Technical complements (3 additional technologies that could complement this invention)
            
            Make all suggestions specific and technical, not general business advice.
            Explain exactly what technical changes would differentiate this patent.
            
            Return as a JSON object with these five properties, each containing an array of strings.
            """
            
            response = self.client.chat.completions.create(
                model=Constants.GPT_MODEL_4,
                messages=[{"role": "user", "content": prompt}]
            )
            
            differentiation = json.loads(response.choices[0].message.content)
            return differentiation
            
        except Exception as e:
            print(f"Error generating differentiation strategies: {e}")
            return {
                "key_competing_patents": [Constants.DEFAULT_ERROR_MESSAGE],
                "technical_differentiation_opportunities": [Constants.DEFAULT_ERROR_MESSAGE],
                "improvement_suggestions": [Constants.DEFAULT_ERROR_MESSAGE],
                "patentability_strategies": [Constants.DEFAULT_ERROR_MESSAGE],
                "technical_complements": [Constants.DEFAULT_ERROR_MESSAGE]
            }
    
    def analyze_market_services(self, idea_description, semantic_keywords=None):
        """시장 서비스 분석"""
        try:
            # 서비스 목록 가져오기
            services = self._get_market_services()
            
            # Analyze services in relation to the user idea
            try:
                prompt = f"""
                Compare the user's patent idea with these example market solutions.
                """
                response = openai.ChatCompletion.create(
                    model=Constants.GPT_MODEL,
                    messages=[
                        {"role": "system", "content": "You are a market analysis expert."},
                        {"role": "user", "content": prompt}
                    ]
                )
                
                analyzed_services = json.loads(response.choices[0].message.content)
                return analyzed_services
                
            except Exception as e:
                print(f"Error analyzing services with GPT: {str(e)}")
                return services
            
        except Exception as e:
            print(f"Error in market services analysis: {str(e)}")
            return []

    def create_powerpoint_presentation(self, idea_description, analysis_results):
        """파워포인트 프레젠테이션 생성"""
        try:
            prs = Presentation()
            
            # Left side: solution description
            left_content = slide.placeholders[1]
            solution_prompt = f"""
            Extract the specific technical solution from this patent idea, formatted as 3-4 bullet points.
            Focus on how the solution works technically, with specific mechanism details.
            """
            
            # Right side: solution image if available
            right_content = slide.placeholders[2]
            if "solution_image" in analysis_results and analysis_results["solution_image"].get("image_url"):
                try:
                    img_response = requests.get(analysis_results["solution_image"]["image_url"])
                    img_stream = io.BytesIO(img_response.content)
                    right_content.insert_picture(img_stream)
                except:
                    right_content.text = "Solution Illustration\n(Image not available)"
                    set_text_size(right_content.text_frame)
            else:
                right_content.text = "Solution Illustration\n(Image not available)"
                set_text_size(right_content.text_frame)
            
            # Add Appendix slide (with differentiation strategies)
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            title.text = "Appendix: Differentiation Strategies"
            set_text_size(title.text_frame, 20, False)
            
            content = slide.placeholders[1]
            
            diff_text = ""
            if differentiation:
                for strategy in differentiation:
                    diff_text += f"• {strategy}\n"
            content.text = diff_text
            set_text_size(content.text_frame)
            
        except Exception as e:
            print(f"Error creating PowerPoint presentation: {str(e)}")
            return None
    
    def _calculate_service_relevance(self, title, description, core_terms, idea_description):
        """Calculate relevance score between market service and patent idea"""
        score = 0
        
        # Calculate term frequency
        for term in core_terms:
            if term.lower() in title.lower():
                score += 2.0
            if term.lower() in description.lower():
                score += 1.0
        
        # Check for product/service indicators
        product_terms = ["product", "solution", "system", "device", "appliance", "technology"]
        for term in product_terms:
            if term.lower() in (title.lower() + " " + description.lower()):
                score += 0.5
        
        # Check for commercial indicators
        commercial_terms = ["buy", "purchase", "order", "subscription", "customer", "price", "cost"]
        for term in commercial_terms:
            if term.lower() in (title.lower() + " " + description.lower()):
                score += 0.5
        
        return score
    
    def run_full_analysis(self, idea_description):
        """전체 특허 분석 실행"""
        try:
            # 특허 검색
            patents = self.search_patents(idea_description)
            
            # 특허 분석
            analyzed_patents = self.analyze_patents(idea_description, patents)
            
            # 차별화 전략 추천
            differentiation = self.recommend_differentiation(idea_description, analyzed_patents)
            
            # 시장 서비스 분석
            market_services = self.analyze_market_services(idea_description)
            
            # 결과 정리
            results = {
                "analyzed_patents": analyzed_patents,
                "differentiation": differentiation,
                "market_services": market_services
            }
            
            return results
            
        except Exception as e:
            print(f"전체 분석 중 오류 발생: {str(e)}")
            traceback.print_exc()
            return None
    
    def get_patents_details(self, patents, max_workers=5):
        """여러 특허의 상세 정보를 병렬로 가져오기"""
        print(f"병렬로 {len(patents)}개의 특허 상세 정보를 가져오는 중...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 각 특허마다 detail 정보 가져오기 작업 제출
            future_to_patent = {
                executor.submit(self.get_patent_details, patent['patent_id']): 
                (i, patent) for i, patent in enumerate(patents)
            }
            
            # 작업이 완료되는 대로 처리
            completed = 0
            for future in concurrent.futures.as_completed(future_to_patent):
                i, patent = future_to_patent[future]
                try:
                    details = future.result()
                    if details:
                        # 원본 특허 정보에 상세 정보 병합
                        patents[i].update(details)
                        
                        # 관련성 점수 계산 (기존 로직 재사용)
                        if hasattr(self, '_calculate_semantic_relevance'):
                            relevance_score = self._calculate_semantic_relevance(
                                patents[i].get('title', ''),
                                patents[i].get('description', ''),
                                getattr(self, 'semantic_keywords', {}),
                                getattr(self, 'idea_description', ''),
                                getattr(self, 'domain_info', {})
                            )
                            patents[i]['relevance_score'] = relevance_score
                    
                    completed += 1
                    print(f"진행 상황: {completed}/{len(patents)} 특허 처리 완료")
                    
                except Exception as e:
                    print(f"특허 {patent['patent_id']} 상세 정보 가져오기 실패: {e}")
                    traceback.print_exc()
        
        return patents
    
    def _calculate_category_match_bonus(self, patent, categories):
        """카테고리 매칭 보너스 점수 계산"""
        bonus = 0.0
        
        # 메인 카테고리 매칭
        if categories["main"].lower() in patent.get('description', '').lower():
            bonus += 0.2
        
        # 서브 카테고리 매칭
        if categories["sub"].lower() in patent.get('description', '').lower():
            bonus += 0.3
        
        # 상세 카테고리 매칭
        for detail in categories["details"]:
            if detail.lower() in patent.get('description', '').lower():
                bonus += 0.1
                
        return min(bonus, 0.5)  # 최대 50% 보너스

    def _extract_category_keywords(self, selected_categories):
        """선택된 카테고리에서 중요 키워드 추출"""
        keywords = {
            "main": [],
            "sub": [],
            "detail": []
        }
        
        # 메인 카테고리 키워드
        main_words = selected_categories["main"].split()
        keywords["main"] = [word.lower() for word in main_words if len(word) > 3]
        
        # 서브 카테고리 키워드
        sub_words = selected_categories["sub"].split()
        keywords["sub"] = [word.lower() for word in sub_words if len(word) > 3]
        
        # 상세 카테고리 키워드
        for detail in selected_categories["details"]:
            detail_words = detail.split()
            keywords["detail"].extend([word.lower() for word in detail_words if len(word) > 3])
        
        return keywords

    def _get_domain_specific_analysis(self, idea_description, main_category):
        """카테고리별 특화된 분석 수행"""
        prompt = f"""
        Analyze this patent idea specifically for the {main_category} domain.
        Consider industry standards, common technical challenges, and innovation opportunities.
        
        Patent idea:
        {idea_description}
        
        Return a JSON object with:
        1. domain_specific_challenges: Array of technical challenges specific to this domain
        2. industry_standards: Key industry standards or regulations to consider
        3. innovation_opportunities: Potential areas for innovation in this domain
        """
        
        try:
            response = self.client.chat.completions.create(
                model=Constants.GPT_MODEL,
                messages=[{"role": "user", "content": prompt}]
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            print(f"도메인 분석 중 오류 발생: {e}")
            return {
                "domain_specific_challenges": [],
                "industry_standards": [],
                "innovation_opportunities": []
            }

    def _optimize_search_query(self, idea_description, category_keywords, domain_focus):
        """카테고리 기반 검색 쿼리 최적화"""
        # 기본 키워드 추출
        base_keywords = self._extract_semantic_keywords(idea_description)
        
        # 카테고리 키워드 통합
        all_keywords = []
        all_keywords.extend(category_keywords["main"])
        all_keywords.extend(category_keywords["sub"][:2])  # 상위 2개만
        all_keywords.extend(category_keywords["detail"][:3])  # 상위 3개만
        
        # 도메인 특화 키워드 추가
        if domain_focus:
            for challenge in domain_focus.get("domain_specific_challenges", [])[:2]:
                all_keywords.extend(challenge.split()[:2])
        
        # 중복 제거 및 정렬
        unique_keywords = list(set(all_keywords))
        
        # 가장 관련성 높은 키워드 선택 (최대 5개)
        top_keywords = sorted(unique_keywords, 
                             key=lambda x: len([k for k in base_keywords["all"] if x in k]), 
                             reverse=True)[:5]
        
        # 검색 쿼리 생성
        search_query = " ".join(top_keywords)
        return search_query

    def analyze_with_categories(self, idea_description, selected_categories):
        """선택된 카테고리를 기반으로 특허 분석 강화"""
        
        # 1. 카테고리별 키워드 추출
        category_keywords = self._extract_category_keywords(selected_categories)
        
        # 2. 도메인 특화 분석
        domain_focus = self._get_domain_specific_analysis(
            idea_description, 
            selected_categories["main"]
        )
        
        # 3. 카테고리별 특허 검색 쿼리 최적화
        search_query = self._optimize_search_query(
            idea_description,
            category_keywords,
            domain_focus
        )
        
        return {
            "category_keywords": category_keywords,
            "domain_focus": domain_focus,
            "optimized_query": search_query
            }

    def generate_concept_image(self, idea_description, image_type="problem"):
        """Generate concept image using DALL-E"""
        try:
            if image_type == "problem":
                prompt = f"Technical patent diagram showing problem: {idea_description}. White background, blue lines, professional patent style with labels and arrows. Simple and clean design suitable for patent documentation."
            else:
                prompt = f"Technical patent diagram showing solution: {idea_description}. White background, blue lines, professional patent style with labels and arrows showing workflow. Simple and clean design suitable for patent documentation."

            print(f"Generating image using DALL-E for: {prompt[:100]}...")
            
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                n=1
            )
            
            return {
                "image_url": response.data[0].url,
                "prompt": prompt
            }
            
        except Exception as e:
            print(f"Error generating image: {str(e)}")
            return None

    @ErrorHandler.handle_api_error
    def analyze_idea(self, idea_description: str) -> Dict:
        """아이디어를 분석하고 결과를 반환합니다."""
        try:
            # 도메인 분류
            self.domain_info = self._classify_patent_domain(idea_description)
            
            # 키워드 추출
            self.semantic_keywords = self._extract_semantic_keywords(idea_description)
            
            # 기술적 특징 분석
            technical_features = self._analyze_technical_features(idea_description)
            
            # 개선 제안 생성
            improvements = self._generate_improvement_suggestions(idea_description)
            
            # 특허성 강화 전략 생성
            strategies = self._generate_patentability_strategies(idea_description)
            
            # 기술적 보완사항 분석
            complements = self._analyze_technical_complements(idea_description)
            
            return {
                'technical_features': technical_features,
                'improvement_suggestions': improvements,
                'patentability_strategies': strategies,
                'technical_complements': complements
            }
            
        except Exception as e:
            logging.error(f"아이디어 분석 중 오류 발생: {str(e)}")
            raise PatentAnalysisError(f"아이디어 분석 중 오류가 발생했습니다: {str(e)}")

    def _analyze_technical_features(self, idea_description: str) -> List[str]:
        """Analyze technical features"""
        prompt = f"""
        Please analyze the key technical features of the following idea:
        {idea_description}
        
        Please provide clear and specific descriptions for each feature.
        """
        
        response = self.client.chat.completions.create(
            model=Constants.GPT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        features = response.choices[0].message.content.split('\n')
        return [f.strip('- ') for f in features if f.strip()]

    def _generate_improvement_suggestions(self, idea_description: str) -> List[str]:
        """Generate improvement suggestions"""
        prompt = f"""
        Please suggest potential improvements for the following idea:
        {idea_description}
        
        Each suggestion should be specific and implementable.
        """
        
        response = self.client.chat.completions.create(
            model=Constants.GPT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        suggestions = response.choices[0].message.content.split('\n')
        return [s.strip('- ') for s in suggestions if s.strip()]

    def _generate_patentability_strategies(self, idea_description: str) -> List[str]:
        """Generate patentability strategies"""
        prompt = f"""
        Please suggest strategies to strengthen the patentability of the following idea:
        {idea_description}
        
        Each strategy should include unique aspects that can be emphasized during patent filing.
        """
        
        response = self.client.chat.completions.create(
            model=Constants.GPT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        strategies = response.choices[0].message.content.split('\n')
        return [s.strip('- ') for s in strategies if s.strip()]

    def _analyze_technical_complements(self, idea_description: str) -> List[str]:
        """Analyze technical complements"""
        prompt = f"""
        Please analyze areas that need technical complementation for the following idea:
        {idea_description}
        
        Each complement should include specific and implementable solutions.
        """
        
        response = self.client.chat.completions.create(
            model=Constants.GPT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        
        complements = response.choices[0].message.content.split('\n')
        return [c.strip('- ') for c in complements if c.strip()]

    def _create_technical_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate technical analysis prompt"""
        return f"""
        Please analyze the technical aspects of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        Inventors: {', '.join(patent_data.get('inventors', []))}
        Assignees: {', '.join(patent_data.get('assignees', []))}
        
        Please include the following aspects in your analysis:
        1. Key Technical Features
        2. Innovation Level
        3. Technical Advantages
        4. Implementation Difficulty
        5. Technical Limitations
        """
    
    def _create_market_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate market analysis prompt"""
        return f"""
        Please analyze the market potential of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Market Size
        2. Target Market
        3. Market Entry Barriers
        4. Profitability
        5. Growth Potential
        """
    
    def _create_competitive_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate competitive analysis prompt"""
        return f"""
        Please analyze the competitive position of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Competitor Landscape
        2. Competitive Advantages
        3. Differentiation Points
        4. Market Share
        5. Competitive Strategy
        """
    
    def _create_risk_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate risk analysis prompt"""
        return f"""
        Please analyze the risks associated with the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Technical Risks
        2. Market Risks
        3. Legal Risks
        4. Operational Risks
        5. Mitigation Strategies
        """
    
    def _create_recommendations_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate recommendations prompt"""
        return f"""
        Please provide recommendations for the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your recommendations:
        1. Technical Improvements
        2. Market Entry Strategy
        3. Competitive Advantage Enhancement
        4. Risk Management Approach
        5. Future Development Direction
        """
    
    def _parse_technical_analysis(self, response: str) -> Dict[str, Any]:
        """Parse technical analysis results"""
        try:
            return {
                "key_features": self._extract_key_points(response, "Key Technical Features"),
                "innovation": self._extract_key_points(response, "Innovation Level"),
                "advantages": self._extract_key_points(response, "Technical Advantages"),
                "difficulty": self._extract_key_points(response, "Implementation Difficulty"),
                "limitations": self._extract_key_points(response, "Technical Limitations")
            }
        except Exception as e:
            logging.error(f"Error parsing technical analysis: {str(e)}")
            return {"error": "Failed to parse technical analysis results"}
    
    def _parse_market_analysis(self, response: str) -> Dict[str, Any]:
        """Parse market analysis results"""
        try:
            return {
                "market_size": self._extract_key_points(response, "Market Size"),
                "target_market": self._extract_key_points(response, "Target Market"),
                "entry_barriers": self._extract_key_points(response, "Market Entry Barriers"),
                "profitability": self._extract_key_points(response, "Profitability"),
                "growth_potential": self._extract_key_points(response, "Growth Potential")
            }
        except Exception as e:
            logging.error(f"Error parsing market analysis: {str(e)}")
            return {"error": "Failed to parse market analysis results"}
    
    def _parse_competitive_analysis(self, response: str) -> Dict[str, Any]:
        """Parse competitive analysis results"""
        try:
            return {
                "competitors": self._extract_key_points(response, "Competitor Landscape"),
                "advantages": self._extract_key_points(response, "Competitive Advantages"),
                "differentiation": self._extract_key_points(response, "Differentiation Points"),
                "market_share": self._extract_key_points(response, "Market Share"),
                "strategy": self._extract_key_points(response, "Competitive Strategy")
            }
        except Exception as e:
            logging.error(f"Error parsing competitive analysis: {str(e)}")
            return {"error": "Failed to parse competitive analysis results"}
    
    def _parse_risk_analysis(self, response: str) -> Dict[str, Any]:
        """Parse risk analysis results"""
        try:
            return {
                "technical_risks": self._extract_key_points(response, "Technical Risks"),
                "market_risks": self._extract_key_points(response, "Market Risks"),
                "legal_risks": self._extract_key_points(response, "Legal Risks"),
                "operational_risks": self._extract_key_points(response, "Operational Risks"),
                "mitigation": self._extract_key_points(response, "Mitigation Strategies")
            }
        except Exception as e:
            logging.error(f"Error parsing risk analysis: {str(e)}")
            return {"error": "Failed to parse risk analysis results"}
    
    def _parse_recommendations(self, response: str) -> Dict[str, Any]:
        """Parse recommendations results"""
        try:
            return {
                "technical_improvements": self._extract_key_points(response, "Technical Improvements"),
                "market_entry": self._extract_key_points(response, "Market Entry Strategy"),
                "competitive_advantage": self._extract_key_points(response, "Competitive Advantage Enhancement"),
                "risk_management": self._extract_key_points(response, "Risk Management Approach"),
                "development_direction": self._extract_key_points(response, "Future Development Direction")
            }
        except Exception as e:
            logging.error(f"Error parsing recommendations: {str(e)}")
            return {"error": "Failed to parse recommendations results"}
    
    def _extract_key_points(self, text: str, section: str) -> List[str]:
        """Extract key points from specific section in text"""
        try:
            # Find section start
            start_idx = text.find(section)
            if start_idx == -1:
                return []
            
            # Find section end
            end_idx = text.find("\n", start_idx)
            if end_idx == -1:
                end_idx = len(text)
            
            # Extract section content
            section_text = text[start_idx:end_idx]
            
            # Extract key points (lines starting with numbers or bullet points)
            points = []
            for line in section_text.split("\n"):
                line = line.strip()
                if line and (line[0].isdigit() or line[0] in ["•", "-", "*"]):
                    points.append(line.lstrip("•-* 123456789."))
            
            return points
        except Exception as e:
            logging.error(f"Error extracting key points: {str(e)}")
            return []

    def analyze_patent_similarity(self, query: str, patent: Dict) -> Dict[str, Any]:
        """특허와 쿼리 간의 유사도를 분석합니다.
        
        Args:
            query: 사용자의 검색 쿼리나 아이디어 설명
            patent: 분석할 특허 데이터
            
        Returns:
            유사도 분석 결과를 포함하는 딕셔너리
        """
        try:
            # 유사도 설명 생성
            similarity_explanation = self._calculate_similarity_explanation(query, patent)
            
            # 시맨틱 유사도 계산
            semantic_score = self._calculate_semantic_relevance(query, patent)
            
            # 키워드 기반 유사도 계산
            keyword_score = self._calculate_keyword_similarity(query, patent)
            
            # 최종 유사도 점수 계산 (0-100 범위)
            # 가중치 조정: 시맨틱 60%, 키워드 40%
            # 추가 페널티: 기술 분야 불일치, 연도 차이 등
            base_score = (semantic_score * 0.6 + keyword_score * 0.4) * 100
            
            # 페널티 계산
            penalties = []
            
            # 기술 분야 불일치 페널티
            if patent.get('cpc_classifications'):
                domain_match = self._check_domain_match(query, patent['cpc_classifications'])
                if not domain_match:
                    penalties.append(0.2)  # 20% 페널티
            
            # 연도 차이 페널티
            if patent.get('publication_date'):
                year_diff = self._calculate_year_difference(patent['publication_date'])
                if year_diff > 10:
                    penalties.append(0.1)  # 10% 페널티
            
            # 페널티 적용
            final_score = base_score
            for penalty in penalties:
                final_score *= (1 - penalty)
            
            return {
                'similarity_explanation': similarity_explanation,
                'semantic_score': semantic_score,
                'keyword_score': keyword_score,
                'relevance_score': min(100, max(0, final_score))
            }
            
        except Exception as e:
            logging.error(f"Error analyzing patent similarity: {str(e)}")
            return {
                'similarity_explanation': "Could not generate similarity explanation due to an error.",
                'semantic_score': 0.5,
                'keyword_score': 0.5,
                'relevance_score': 50
            }

    def _check_domain_match(self, query: str, cpc_classifications: List[str]) -> bool:
        """쿼리와 특허의 기술 분야가 일치하는지 확인합니다."""
        try:
            # 쿼리에서 주요 기술 분야 추출
            domain_prompt = f"""
            Extract the main technical domain from this query:
            {query}
            
            Return only the main technical domain (e.g., "AI", "Robotics", "Biotechnology").
            """
            
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": domain_prompt}],
                temperature=0.3,
                max_tokens=20
            )
            
            query_domain = response.choices[0].message.content.strip().lower()
            
            # CPC 분류와 비교
            for cpc in cpc_classifications:
                if query_domain in cpc.lower():
                    return True
            
            return False
            
        except Exception as e:
            logging.error(f"Error checking domain match: {str(e)}")
            return False

    def _calculate_year_difference(self, publication_date: str) -> int:
        """현재 연도와 특허 출원 연도의 차이를 계산합니다."""
        try:
            current_year = datetime.now().year
            patent_year = int(publication_date.split('-')[0])
            return current_year - patent_year
        except:
            return 0

    def _calculate_similarity_explanation(self, query: str, patent: Dict) -> str:
        """특허와 쿼리 간의 유사성을 설명하는 상세한 분석을 생성합니다."""
        try:
            # 특허 정보 추출
            patent_title = patent.get('title', '')
            patent_abstract = patent.get('abstract', '')
            patent_description = patent.get('description', '')
            patent_claims = patent.get('claims', [])
            
            # 클레임 텍스트 포함 (최대 3개)
            claims_text = ""
            if patent_claims and len(patent_claims) > 0:
                claims_text = "Claims:\n" + "\n".join([f"- {claim}" for claim in patent_claims[:3]])
            
            # 분석 프롬프트 생성
            prompt = f"""
            Compare this user's patent idea with an existing patent and explain in detail WHY they are similar or different.
            
            USER'S IDEA:
            {query}
            
            EXISTING PATENT:
            Title: {patent_title}
            Abstract: {patent_abstract}
            Description: {patent_description}
            {claims_text}
            
            Provide a detailed technical explanation (200+ words) covering:
            1. Core technological similarities and differences
            2. Implementation approaches comparison
            3. Potential overlapping claims
            4. Unique aspects of each
            5. How the existing patent might impact the patentability of the user's idea
            
            Format the response as a detailed technical analysis focused on helping the user understand the relevance of this patent to their invention.
            """
            
            # OpenAI API 호출
            response = self.client.chat.completions.create(
                model="gpt-4",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.5,
                max_tokens=1000
            )
            
            return response.choices[0].message.content.strip()
            
        except Exception as e:
            logging.error(f"Error generating similarity explanation: {str(e)}")
            return "Could not generate detailed similarity explanation due to an error."

    def _calculate_semantic_relevance(self, query: str, patent: Dict) -> float:
        """시맨틱 유사도 점수를 계산합니다 (0-1 범위)"""
        try:
            # 특허 텍스트 준비
            patent_text = f"{patent.get('title', '')} {patent.get('abstract', '')} {patent.get('description', '')}"
            
            # 임베딩 생성을 위한 프롬프트
            prompt = f"""
            Rate the semantic relevance between these two texts on a scale of 0 to 1:
            
            Text 1: {query}
            
            Text 2: {patent_text}
            
            Return only a number between 0 and 1, where 1 means highly relevant and 0 means not relevant at all.
            """
            
            # OpenAI API 호출
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=10
            )
            
            # 응답에서 숫자 추출
            score_text = response.choices[0].message.content.strip()
            score = float(re.search(r'0*\.?\d+', score_text).group())
            return min(1.0, max(0.0, score))
            
        except Exception as e:
            logging.error(f"Error calculating semantic relevance: {str(e)}")
            return 0.5

    def _calculate_keyword_similarity(self, query: str, patent: Dict) -> float:
        """키워드 기반 유사도 점수를 계산합니다 (0-1 범위)"""
        try:
            # 쿼리에서 키워드 추출
            keywords = [word.lower() for word in query.split() if len(word) > 3]
            
            # 특허 텍스트 정규화
            title = patent.get('title', '').lower()
            abstract = patent.get('abstract', '').lower()
            description = patent.get('description', '').lower()
            
            score = 0.0
            
            # 제목에서 키워드 매칭 (높은 가중치)
            for keyword in keywords:
                if keyword in title:
                    score += 0.15
                    # 정확한 단어 매칭 확인
                    if re.search(r'\b' + re.escape(keyword) + r'\b', title):
                        score += 0.05
            
            # 초록에서 키워드 매칭 (중간 가중치)
            for keyword in keywords:
                if keyword in abstract:
                    score += 0.1
                    if re.search(r'\b' + re.escape(keyword) + r'\b', abstract):
                        score += 0.03
            
            # 설명에서 키워드 매칭 (낮은 가중치)
            for keyword in keywords:
                if keyword in description:
                    score += 0.05
                    if re.search(r'\b' + re.escape(keyword) + r'\b', description):
                        score += 0.02
            
            return min(1.0, score)
            
        except Exception as e:
            logging.error(f"Error calculating keyword similarity: {str(e)}")
            return 0.5

class PatentAnalyzer:
    """Class for patent analysis"""
    
    def __init__(self, openai_api_key: str):
        """Initialize patent analyzer"""
        self.openai_client = OpenAI(api_key=openai_api_key)
        self.cache_manager = CacheManager()
        self.error_handler = ErrorHandler()
    
    @ErrorHandler.handle_api_error
    def analyze_patent(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze patent data"""
        cache_key = f"patent_analysis_{patent_data.get('patent_id', 'unknown')}"
        cached_result = self.cache_manager.get_cached_data(cache_key)
        
        if cached_result:
            return cached_result
            
        analysis_result = self._perform_analysis(patent_data)
        self.cache_manager.save_to_cache(cache_key, analysis_result)
        return analysis_result
    
    def _perform_analysis(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Perform actual analysis"""
        try:
            return {
                "technical_analysis": self._analyze_technical_aspects(patent_data),
                "market_analysis": self._analyze_market_potential(patent_data),
                "competitive_analysis": self._analyze_competitive_position(patent_data),
                "risk_analysis": self._analyze_risks(patent_data),
                "recommendations": self._generate_recommendations(patent_data)
            }
        except Exception as e:
            logging.error(f"Error during analysis: {str(e)}")
            raise PatentAnalysisError(f"Error occurred during patent analysis: {str(e)}")
    
    def _analyze_technical_aspects(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze technical aspects"""
        prompt = self._create_technical_analysis_prompt(patent_data)
        response = self._get_gpt_response(prompt)
        return self._parse_technical_analysis(response)
    
    def _analyze_market_potential(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze market potential"""
        prompt = self._create_market_analysis_prompt(patent_data)
        response = self._get_gpt_response(prompt)
        return self._parse_market_analysis(response)
    
    def _analyze_competitive_position(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze competitive position"""
        prompt = self._create_competitive_analysis_prompt(patent_data)
        response = self._get_gpt_response(prompt)
        return self._parse_competitive_analysis(response)
    
    def _analyze_risks(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze risks"""
        prompt = self._create_risk_analysis_prompt(patent_data)
        response = self._get_gpt_response(prompt)
        return self._parse_risk_analysis(response)
    
    def _generate_recommendations(self, patent_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate recommendations"""
        prompt = self._create_recommendations_prompt(patent_data)
        response = self._get_gpt_response(prompt)
        return self._parse_recommendations(response)
    
    def _get_gpt_response(self, prompt: str) -> str:
        """Get response from GPT model"""
        try:
            response = self.openai_client.chat.completions.create(
                model=Constants.GPT_MODEL,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.choices[0].message.content
        except Exception as e:
            logging.error(f"Error during GPT API call: {str(e)}")
            raise PatentAnalysisError(f"Error occurred during GPT API call: {str(e)}")
    
    def _create_technical_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate technical analysis prompt"""
        return f"""
        Please analyze the technical aspects of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        Inventors: {', '.join(patent_data.get('inventors', []))}
        Assignees: {', '.join(patent_data.get('assignees', []))}
        
        Please include the following aspects in your analysis:
        1. Key Technical Features
        2. Innovation Level
        3. Technical Advantages
        4. Implementation Difficulty
        5. Technical Limitations
        """
    
    def _create_market_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate market analysis prompt"""
        return f"""
        Please analyze the market potential of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Market Size
        2. Target Market
        3. Market Entry Barriers
        4. Profitability
        5. Growth Potential
        """
    
    def _create_competitive_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate competitive analysis prompt"""
        return f"""
        Please analyze the competitive position of the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Competitor Landscape
        2. Competitive Advantages
        3. Differentiation Points
        4. Market Share
        5. Competitive Strategy
        """
    
    def _create_risk_analysis_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate risk analysis prompt"""
        return f"""
        Please analyze the risks associated with the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your analysis:
        1. Technical Risks
        2. Market Risks
        3. Legal Risks
        4. Operational Risks
        5. Mitigation Strategies
        """
    
    def _create_recommendations_prompt(self, patent_data: Dict[str, Any]) -> str:
        """Generate recommendations prompt"""
        return f"""
        Please provide recommendations for the following patent:
        Title: {patent_data.get('title', 'N/A')}
        Abstract: {patent_data.get('abstract', 'N/A')}
        
        Please include the following aspects in your recommendations:
        1. Technical Improvements
        2. Market Entry Strategy
        3. Competitive Advantage Enhancement
        4. Risk Management Approach
        5. Future Development Direction
        """
    
    def _parse_technical_analysis(self, response: str) -> Dict[str, Any]:
        """Parse technical analysis results"""
        try:
            return {
                "key_features": self._extract_key_points(response, "Key Technical Features"),
                "innovation": self._extract_key_points(response, "Innovation Level"),
                "advantages": self._extract_key_points(response, "Technical Advantages"),
                "difficulty": self._extract_key_points(response, "Implementation Difficulty"),
                "limitations": self._extract_key_points(response, "Technical Limitations")
            }
        except Exception as e:
            logging.error(f"Error parsing technical analysis: {str(e)}")
            return {"error": "Failed to parse technical analysis results"}
    
    def _parse_market_analysis(self, response: str) -> Dict[str, Any]:
        """Parse market analysis results"""
        try:
            return {
                "market_size": self._extract_key_points(response, "Market Size"),
                "target_market": self._extract_key_points(response, "Target Market"),
                "entry_barriers": self._extract_key_points(response, "Market Entry Barriers"),
                "profitability": self._extract_key_points(response, "Profitability"),
                "growth_potential": self._extract_key_points(response, "Growth Potential")
            }
        except Exception as e:
            logging.error(f"Error parsing market analysis: {str(e)}")
            return {"error": "Failed to parse market analysis results"}
    
    def _parse_competitive_analysis(self, response: str) -> Dict[str, Any]:
        """Parse competitive analysis results"""
        try:
            return {
                "competitors": self._extract_key_points(response, "Competitor Landscape"),
                "advantages": self._extract_key_points(response, "Competitive Advantages"),
                "differentiation": self._extract_key_points(response, "Differentiation Points"),
                "market_share": self._extract_key_points(response, "Market Share"),
                "strategy": self._extract_key_points(response, "Competitive Strategy")
            }
        except Exception as e:
            logging.error(f"Error parsing competitive analysis: {str(e)}")
            return {"error": "Failed to parse competitive analysis results"}
    
    def _parse_risk_analysis(self, response: str) -> Dict[str, Any]:
        """Parse risk analysis results"""
        try:
            return {
                "technical_risks": self._extract_key_points(response, "Technical Risks"),
                "market_risks": self._extract_key_points(response, "Market Risks"),
                "legal_risks": self._extract_key_points(response, "Legal Risks"),
                "operational_risks": self._extract_key_points(response, "Operational Risks"),
                "mitigation": self._extract_key_points(response, "Mitigation Strategies")
            }
        except Exception as e:
            logging.error(f"Error parsing risk analysis: {str(e)}")
            return {"error": "Failed to parse risk analysis results"}
    
    def _parse_recommendations(self, response: str) -> Dict[str, Any]:
        """Parse recommendations results"""
        try:
            return {
                "technical_improvements": self._extract_key_points(response, "Technical Improvements"),
                "market_entry": self._extract_key_points(response, "Market Entry Strategy"),
                "competitive_advantage": self._extract_key_points(response, "Competitive Advantage Enhancement"),
                "risk_management": self._extract_key_points(response, "Risk Management Approach"),
                "development_direction": self._extract_key_points(response, "Future Development Direction")
            }
        except Exception as e:
            logging.error(f"Error parsing recommendations: {str(e)}")
            return {"error": "Failed to parse recommendations results"}
    
    def _extract_key_points(self, text: str, section: str) -> List[str]:
        """Extract key points from specific section in text"""
        try:
            # Find section start
            start_idx = text.find(section)
            if start_idx == -1:
                return []
            
            # Find section end
            end_idx = text.find("\n", start_idx)
            if end_idx == -1:
                end_idx = len(text)
            
            # Extract section content
            section_text = text[start_idx:end_idx]
            
            # Extract key points (lines starting with numbers or bullet points)
            points = []
            for line in section_text.split("\n"):
                line = line.strip()
                if line and (line[0].isdigit() or line[0] in ["•", "-", "*"]):
                    points.append(line.lstrip("•-* 123456789."))
            
            return points
        except Exception as e:
            logging.error(f"Error extracting key points: {str(e)}")
            return []

class PatentVisualizer:
    """Class for patent data visualization"""
    
    def __init__(self):
        """Initialize the visualizer"""
        self.error_handler = ErrorHandler()
    
    @ErrorHandler.handle_file_operation
    def create_presentation(self, patent_data: Dict[str, Any], analysis_result: Dict[str, Any], output_path: str) -> None:
        """Create a PowerPoint presentation from patent analysis results"""
        prs = Presentation()
        
        # Add title slide
        self._add_title_slide(prs, patent_data)
        
        # Add technical analysis slide
        self._add_technical_analysis_slide(prs, analysis_result)
        
        # Add market analysis slide
        self._add_market_analysis_slide(prs, analysis_result)
        
        # Add competitive analysis slide
        self._add_competitive_analysis_slide(prs, analysis_result)
        
        # Add risk analysis slide
        self._add_risk_analysis_slide(prs, analysis_result)
        
        # Add recommendations slide
        self._add_recommendations_slide(prs, analysis_result)
        
        # Save presentation
        prs.save(output_path)
    
    def _add_title_slide(self, prs: Presentation, patent_data: Dict[str, Any]) -> None:
        """Add title slide to presentation"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Patent Analysis: {patent_data.get('title', 'N/A')}"
        subtitle.text = f"Patent ID: {patent_data.get('patent_id', 'N/A')}"
    
    def _add_technical_analysis_slide(self, prs: Presentation, analysis_result: Dict[str, Any]) -> None:
        """Add technical analysis slide to presentation"""
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Technical Analysis"
        content.text = self._format_analysis_content(analysis_result.get("technical_analysis", {}))
    
    def _add_market_analysis_slide(self, prs: Presentation, analysis_result: Dict[str, Any]) -> None:
        """Add market analysis slide to presentation"""
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Market Analysis"
        content.text = self._format_analysis_content(analysis_result.get("market_analysis", {}))
    
    def _add_competitive_analysis_slide(self, prs: Presentation, analysis_result: Dict[str, Any]) -> None:
        """Add competitive analysis slide to presentation"""
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Competitive Analysis"
        content.text = self._format_analysis_content(analysis_result.get("competitive_analysis", {}))
    
    def _add_risk_analysis_slide(self, prs: Presentation, analysis_result: Dict[str, Any]) -> None:
        """Add risk analysis slide to presentation"""
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Risk Analysis"
        content.text = self._format_analysis_content(analysis_result.get("risk_analysis", {}))
    
    def _add_recommendations_slide(self, prs: Presentation, analysis_result: Dict[str, Any]) -> None:
        """Add recommendations slide to presentation"""
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Recommendations"
        content.text = self._format_analysis_content(analysis_result.get("recommendations", {}))
    
    def _format_analysis_content(self, analysis_data: Dict[str, Any]) -> str:
        """Format analysis data for presentation"""
        formatted_text = ""
        
        for key, value in analysis_data.items():
            if isinstance(value, list):
                formatted_text += f"\n{key.replace('_', ' ').title()}:\n"
                for item in value:
                    formatted_text += f"• {item}\n"
                formatted_text += "\n"
        
        return formatted_text
    
    @ErrorHandler.handle_file_operation
    def create_flowchart(self, patent_data: Dict[str, Any], output_path: str) -> None:
        """Create a flowchart of the patent process"""
        dot = graphviz.Digraph(comment='Patent Process Flowchart')
        dot.attr(rankdir='TB')
        
        # Add nodes
        dot.node('A', 'Patent Search')
        dot.node('B', 'Technical Analysis')
        dot.node('C', 'Market Analysis')
        dot.node('D', 'Competitive Analysis')
        dot.node('E', 'Risk Analysis')
        dot.node('F', 'Recommendations')
        
        # Add edges
        dot.edge('A', 'B')
        dot.edge('B', 'C')
        dot.edge('C', 'D')
        dot.edge('D', 'E')
        dot.edge('E', 'F')
        
        # Save flowchart
        dot.render(output_path, cleanup=True)
    
    @ErrorHandler.handle_file_operation
    def create_radar_chart(self, analysis_result: Dict[str, Any], output_path: str) -> None:
        """Create a radar chart of patent analysis results"""
        # Calculate scores for each category
        technical_score = self._calculate_technical_score(analysis_result)
        market_score = self._calculate_market_score(analysis_result)
        competitive_score = self._calculate_competitive_score(analysis_result)
        risk_score = self._calculate_risk_score(analysis_result)
        improvement_score = self._calculate_improvement_score(analysis_result)
        
        # Create radar chart
        fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
        
        # Data for radar chart
        categories = ['Technical', 'Market', 'Competitive', 'Risk', 'Improvement']
        scores = [technical_score, market_score, competitive_score, risk_score, improvement_score]
        
        # Plot radar chart
        angles = np.linspace(0, 2*np.pi, len(categories), endpoint=False)
        scores = np.concatenate((scores, [scores[0]]))  # Close the plot
        angles = np.concatenate((angles, [angles[0]]))  # Close the plot
        
        ax.plot(angles, scores)
        ax.fill(angles, scores, alpha=0.25)
        
        # Set labels
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories)
        
        # Set title
        plt.title('Patent Analysis Radar Chart')
        
        # Save chart
        plt.savefig(output_path)
        plt.close()
    
    def _calculate_technical_score(self, analysis_result: Dict[str, Any]) -> float:
        """Calculate technical analysis score"""
        technical_data = analysis_result.get("technical_analysis", {})
        if not technical_data:
            return 0.0
            
        score = 0.0
        weights = {
            "key_features": 0.3,
            "innovation": 0.2,
            "advantages": 0.2,
            "difficulty": 0.15,
            "limitations": 0.15
        }
        
        for key, weight in weights.items():
            if key in technical_data and isinstance(technical_data[key], list):
                score += len(technical_data[key]) * weight
        
        return min(score, 1.0)
    
    def _calculate_market_score(self, analysis_result: Dict[str, Any]) -> float:
        """Calculate market analysis score"""
        market_data = analysis_result.get("market_analysis", {})
        if not market_data:
            return 0.0
            
        score = 0.0
        weights = {
            "market_size": 0.2,
            "target_market": 0.2,
            "entry_barriers": 0.2,
            "profitability": 0.2,
            "growth_potential": 0.2
        }
        
        for key, weight in weights.items():
            if key in market_data and isinstance(market_data[key], list):
                score += len(market_data[key]) * weight
        
        return min(score, 1.0)
    
    def _calculate_competitive_score(self, analysis_result: Dict[str, Any]) -> float:
        """Calculate competitive analysis score"""
        competitive_data = analysis_result.get("competitive_analysis", {})
        if not competitive_data:
            return 0.0
            
        score = 0.0
        weights = {
            "competitors": 0.2,
            "advantages": 0.2,
            "differentiation": 0.2,
            "market_share": 0.2,
            "strategy": 0.2
        }
        
        for key, weight in weights.items():
            if key in competitive_data and isinstance(competitive_data[key], list):
                score += len(competitive_data[key]) * weight
        
        return min(score, 1.0)
    
    def _calculate_risk_score(self, analysis_result: Dict[str, Any]) -> float:
        """Calculate risk analysis score"""
        risk_data = analysis_result.get("risk_analysis", {})
        if not risk_data:
            return 0.0
            
        score = 0.0
        weights = {
            "technical_risks": 0.2,
            "market_risks": 0.2,
            "legal_risks": 0.2,
            "operational_risks": 0.2,
            "mitigation": 0.2
        }
        
        for key, weight in weights.items():
            if key in risk_data and isinstance(risk_data[key], list):
                score += len(risk_data[key]) * weight
        
        return min(score, 1.0)
    
    def _calculate_improvement_score(self, analysis_result: Dict[str, Any]) -> float:
        """Calculate improvement recommendations score"""
        recommendations = analysis_result.get("recommendations", {})
        if not recommendations:
            return 0.0
            
        score = 0.0
        weights = {
            "technical_improvements": 0.2,
            "market_entry": 0.2,
            "competitive_advantage": 0.2,
            "risk_management": 0.2,
            "development_direction": 0.2
        }
        
        for key, weight in weights.items():
            if key in recommendations and isinstance(recommendations[key], list):
                score += len(recommendations[key]) * weight
        
        return min(score, 1.0)
    
# patent.py 파일에 추가할 클래스 정의
class EnhancedPPTGenerator:
    """Enhanced PPT generator for patent reports"""
    
    def __init__(self, api_key=None):
        """Initialize PPT generator with API key"""
        self.api_key = api_key or os.getenv('OPENAI_API_KEY')
        if not self.api_key:
            raise ConfigurationError("OpenAI API key is required")
        self.client = OpenAI(api_key=self.api_key)
        self.font_size = Pt(15)
    
    def create_presentation(self, idea_description: str, patent_data: List[Dict], analysis_result: Dict, output_path: str, include_similarity=True, include_appendix=True) -> str:
        """Create a presentation with the specified content"""
        prs = Presentation()
        
        # Set slide width and height (16:9 aspect ratio)
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Add slides in order
        self._add_title_slide(prs, idea_description)
        self._add_executive_summary_slide(prs, idea_description, analysis_result)
        self._add_problem_slide(prs, idea_description, analysis_result)
        self._add_solution_slide(prs, idea_description, analysis_result)
        self._add_solution_flowchart_slide(prs, idea_description, analysis_result)
        self._add_benefits_slide(prs, analysis_result)
        self._add_prior_arts_slide(prs, patent_data)
        
        # Add similarity analysis if requested
        if include_similarity and patent_data:
            self._add_similarity_analysis_slide(prs, idea_description, patent_data)
        
        # Add appendix if there's additional technical information
        if include_appendix and analysis_result.get('technical_complements'):
            self._add_appendix_slide(prs, analysis_result)
        
        # Save presentation
        prs.save(output_path)
        return output_path
    
    def _set_text_size(self, text_frame):
        """Set text size for all paragraphs in a text frame"""
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = self.font_size
            for run in paragraph.runs:
                run.font.size = self.font_size
    
    def _add_title_slide(self, prs, idea_description):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # 제목 생성 (아이디어 처음 30자 + "...")
        title_text = idea_description[:30] + "..." if len(idea_description) > 30 else idea_description
        title.text = f"TITLE: {title_text}"
        subtitle.text = f"Date: {datetime.now().strftime('%Y-%m-%d')}"
        
        # 글씨 크기 설정
        self._set_text_size(title.text_frame)
        self._set_text_size(subtitle.text_frame)
    
    def _add_executive_summary_slide(self, prs, idea_description, analysis_result):
        """Add executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "EXECUTIVE SUMMARY"
        
        summary_text = f"Idea Overview:\n{idea_description}\n\n"
        if analysis_result.get('technical_features'):
            summary_text += "Key Technical Features:\n"
            for feature in analysis_result['technical_features'][:3]:
                summary_text += f"• {feature}\n"
        
        content.text = summary_text
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_problem_slide(self, prs, idea_description, analysis_result):
        """Add problem description slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "PROBLEM"
        
        problem_text = "Current Challenges:\n\n"
        if analysis_result.get('problem_statement'):
            problem_text += analysis_result['problem_statement']
        else:
            problem_text += idea_description
        
        content.text = problem_text
        
        # Add problem image if available
        try:
            problem_image = self._generate_problem_image(idea_description)
            if problem_image:
                left = Inches(1)
                top = Inches(3)
                width = Inches(8)
                height = Inches(4.5)
                slide.shapes.add_picture(problem_image, left, top, width=width, height=height)
        except Exception as e:
            logging.error(f"Error adding problem image: {str(e)}")
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_solution_slide(self, prs, idea_description, analysis_result):
        """Add solution slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "SOLUTION"
        
        solution_text = "Proposed Solution:\n\n"
        if analysis_result.get('technical_features'):
            for feature in analysis_result['technical_features']:
                solution_text += f"• {feature}\n"
        
        content.text = solution_text
        
        # Add solution image if available
        try:
            solution_image = self._generate_solution_image(idea_description)
            if solution_image:
                left = Inches(1)
                top = Inches(3)
                width = Inches(8)
                height = Inches(4.5)
                slide.shapes.add_picture(solution_image, left, top, width=width, height=height)
        except Exception as e:
            logging.error(f"Error adding solution image: {str(e)}")
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_solution_flowchart_slide(self, prs, idea_description, analysis_result):
        """Add solution flowchart slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "SOLUTION – FLOW CHART"
        content.text = "Process Flow Diagram"
        
        # Add flowchart
        try:
            flowchart_image = self._generate_flowchart(idea_description, analysis_result)
            if flowchart_image:
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(5.5)
                slide.shapes.add_picture(flowchart_image, left, top, width=width, height=height)
        except Exception as e:
            logging.error(f"Error adding flowchart: {str(e)}")
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_benefits_slide(self, prs, analysis_result):
        """Add benefits slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "BENEFITS"
        
        benefits_text = "1. Benefits for Consumers:\n"
        if analysis_result.get('improvement_suggestions'):
            for imp in analysis_result['improvement_suggestions'][:3]:
                benefits_text += f"• {imp}\n"
        
        benefits_text += "\n2. Benefits for GEA:\n"
        if analysis_result.get('patentability_strategies'):
            for strategy in analysis_result['patentability_strategies'][:3]:
                benefits_text += f"• {strategy}\n"
        
        content.text = benefits_text
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_prior_arts_slide(self, prs, patent_data):
        """Add prior arts slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "PRIOR ARTS"
        
        prior_arts_text = "Similar Patents:\n\n"
        for i, patent in enumerate(patent_data[:5], 1):
            prior_arts_text += f"{i}. {patent.get('title', 'Unknown Title')}\n"
            prior_arts_text += f"   Patent ID: {patent.get('patent_id', 'N/A')}\n"
            prior_arts_text += f"   Assignee: {patent.get('assignee', 'N/A')}\n\n"
        
        content.text = prior_arts_text
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _add_appendix_slide(self, prs, analysis_result):
        """Add appendix slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "APPENDIX"
        
        appendix_text = "Additional Technical Details:\n\n"
        if analysis_result.get('technical_complements'):
            for comp in analysis_result['technical_complements']:
                appendix_text += f"• {comp}\n"
        
        content.text = appendix_text
        
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)
    
    def _generate_problem_image(self, idea_description):
        """Generate image for problem description"""
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=f"Technical patent diagram showing problem: {idea_description}. White background, blue lines, professional patent style with labels and arrows. Simple and clean design suitable for patent documentation.",
                size="1024x1024",
                quality="standard",
                n=1
            )
            
            img_response = requests.get(response.data[0].url)
            return BytesIO(img_response.content)
        except Exception as e:
            logging.error(f"Error generating problem image: {str(e)}")
            return None
    
    def _generate_solution_image(self, idea_description):
        """Generate image for solution description"""
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=f"Technical patent diagram showing solution: {idea_description}. White background, blue lines, professional patent style with labels and arrows showing workflow. Simple and clean design suitable for patent documentation.",
                size="1024x1024",
                quality="standard",
                n=1
            )
            
            img_response = requests.get(response.data[0].url)
            return BytesIO(img_response.content)
        except Exception as e:
            logging.error(f"Error generating solution image: {str(e)}")
            return None

    def _generate_flowchart(self, idea_description, analysis_result):
        """Generate a detailed flowchart using Graphviz"""
        try:
            # Create a new Digraph object
            dot = graphviz.Digraph(comment='Solution Process Flowchart')
            dot.attr(rankdir='TB')
            
            # Set global graph attributes for better appearance
            dot.attr('graph', 
                    fontname='Arial',
                    splines='ortho',
                    nodesep='0.5',
                    ranksep='0.5',
                    bgcolor='white')
            
            # Set default node attributes
            dot.attr('node',
                    shape='box',
                    style='rounded,filled',
                    fillcolor='lightblue',
                    fontname='Arial',
                    fontsize='12',
                    margin='0.2',
                    height='0.6')
            
            # Set default edge attributes
            dot.attr('edge',
                    fontname='Arial',
                    fontsize='10',
                    color='darkblue',
                    arrowsize='0.8')

            # Extract features and improvements from analysis result
            features = analysis_result.get('technical_features', [])
            improvements = analysis_result.get('improvement_suggestions', [])

            # Create nodes based on analysis results
            if features and improvements:
                # Add start node
                dot.node('start', 'Start\nProblem Recognition', shape='oval')
                
                # Add problem identification
                dot.node('problem', self._wrap_text(features[0], 20))
                dot.edge('start', 'problem')
                
                # Add core technical features
                for i, feature in enumerate(features[1:4], 1):
                    node_id = f'tech_{i}'
                    dot.node(node_id, self._wrap_text(feature, 20))
                    if i == 1:
                        dot.edge('problem', node_id)
                    else:
                        dot.edge(f'tech_{i-1}', node_id)
                
                # Add improvement suggestions
                for i, imp in enumerate(improvements[:3], 1):
                    node_id = f'imp_{i}'
                    dot.node(node_id, self._wrap_text(imp, 20))
                    if i == 1:
                        dot.edge('tech_3', node_id)
                    else:
                        dot.edge(f'imp_{i-1}', node_id)
                
                # Add end node
                dot.node('end', 'End\nFinal Solution', shape='oval')
                dot.edge('imp_3', 'end')
            else:
                # Create a basic flowchart if no analysis results
                nodes = [
                    ('start', 'Start', 'oval'),
                    ('prob_id', 'Problem Identification', 'box'),
                    ('analysis', 'Technical Analysis', 'box'),
                    ('solution', 'Solution Development', 'box'),
                    ('implement', 'Implementation Plan', 'box'),
                    ('validate', 'Validation & Testing', 'box'),
                    ('optimize', 'Optimization', 'box'),
                    ('end', 'End', 'oval')
                ]
                
                # Add nodes
                for node_id, label, shape in nodes:
                    dot.node(node_id, label, shape=shape)
                
                # Add edges
                for i in range(len(nodes)-1):
                    dot.edge(nodes[i][0], nodes[i+1][0])

            # Save to a temporary file
            temp_dir = tempfile.gettempdir()
            temp_path = os.path.join(temp_dir, f'flowchart_{uuid.uuid4()}')
            dot.render(temp_path, format='png', cleanup=True)
            
            # Read the generated image
            with open(f'{temp_path}.png', 'rb') as f:
                return BytesIO(f.read())
                
        except Exception as e:
            logging.error(f"Error generating flowchart: {str(e)}")
            return None
    
    def _wrap_text(self, text, width):
        """Wrap text to fit in nodes"""
        return '\n'.join(textwrap.wrap(text, width=width))

    def _add_similarity_analysis_slide(self, prs, idea_description, patent_data):
        """Add a slide showing similarity analysis between the idea and patents"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "SIMILARITY ANALYSIS"
        
        # Find the top 3 most similar patents based on relevance score
        sorted_patents = sorted(
            [p for p in patent_data if 'relevance_score' in p],
            key=lambda x: x.get('relevance_score', 0),
            reverse=True
        )[:3]
        
        similarity_text = "Comparison with Most Similar Patents:\n\n"
        
        for i, patent in enumerate(sorted_patents, 1):
            similarity_text += f"{i}. {patent.get('title', 'Unknown Patent')}\n"
            similarity_text += f"   Relevance Score: {patent.get('relevance_score', 0):.1f}%\n"
            
            # Add key similarities if available
            if 'analysis' in patent and 'similarities' in patent['analysis']:
                similarity_text += "   Key Similarities:\n"
                for sim in patent['analysis']['similarities'][:2]:  # Top 2 similarities
                    similarity_text += f"   - {sim}\n"
            
            # Add key differences if available
            if 'analysis' in patent and 'differences' in patent['analysis']:
                similarity_text += "   Key Differences:\n"
                for diff in patent['analysis']['differences'][:2]:  # Top 2 differences
                    similarity_text += f"   - {diff}\n"
            
            # Add similarity explanation if available
            if 'similarity_explanation' in patent:
                explanation = patent['similarity_explanation']
                # Truncate explanation if too long
                if len(explanation) > 300:
                    explanation = explanation[:297] + "..."
                similarity_text += f"   Summary: {explanation}\n"
            
            similarity_text += "\n"
        
        # Add visualization if available
        try:
            # Create similarity visualization
            fig = plt.figure(figsize=(8, 4))
            scores = [p.get('relevance_score', 0) for p in sorted_patents]
            titles = [p.get('title', 'Unknown')[:30] + "..." for p in sorted_patents]
            
            plt.barh(titles, scores, color=['#2ecc71' if s >= 75 else '#f1c40f' if s >= 50 else '#e74c3c' for s in scores])
            plt.xlabel('Relevance Score (%)')
            plt.title('Patent Similarity Comparison')
            
            # Save to BytesIO
            img_stream = BytesIO()
            plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
            img_stream.seek(0)
            
            # Add image to slide
            left = Inches(1)
            top = Inches(5)
            width = Inches(8)
            slide.shapes.add_picture(img_stream, left, top, width=width)
            plt.close()
            
        except Exception as e:
            logging.error(f"Error creating similarity visualization: {str(e)}")
        
        content.text = similarity_text
        
        # Set text formatting
        self._set_text_size(title.text_frame)
        self._set_text_size(content.text_frame)

def main() -> None:
    """Main function"""
    st.set_page_config(page_title="Patent Analysis System", page_icon="📝")
    
    # Sidebar settings
    st.sidebar.title("Settings")
    openai_api_key = st.sidebar.text_input("OpenAI API Key", type="password")
    if not openai_api_key:
        st.warning("Please enter your OpenAI API key in the sidebar.")
        return
    
    # Patent search section
    st.title("Patent Analysis System")
    search_query = st.text_input("Enter search query for patent search")
    
    if st.button("Search Patents"):
        if not search_query:
            st.error("Please enter a search query.")
            return
        
        try:
            # Initialize patent scraper and analyzer
            scraper = PatentScraper()
            analyzer = PatentAnalyzer(openai_api_key)
            visualizer = PatentVisualizer()
            
            # Search patents
            with st.spinner("Searching patents..."):
                patents = scraper.search_patents(search_query)
            
            if not patents:
                st.warning("No patents found.")
                return
            
            # Display patents and let user select one
            selected_patent = st.selectbox(
                "Select a patent to analyze",
                options=patents,
                format_func=lambda x: f"{x['title']} (ID: {x['patent_id']})"
            )
            
            if st.button("Analyze Patent"):
                # Analyze patent
                with st.spinner("Analyzing patent..."):
                    analysis_result = analyzer.analyze_patent(selected_patent)
                
                # Display analysis results
                st.header("Analysis Results")
                
                # Technical Analysis
                st.subheader("Technical Analysis")
                technical = analysis_result.get("technical_analysis", {})
                for key, value in technical.items():
                    if isinstance(value, list) and value:
                        st.write(f"**{key.replace('_', ' ').title()}:**")
                        for item in value:
                            st.write(f"- {item}")
                
                # Market Analysis
                st.subheader("Market Analysis")
                market = analysis_result.get("market_analysis", {})
                for key, value in market.items():
                    if isinstance(value, list) and value:
                        st.write(f"**{key.replace('_', ' ').title()}:**")
                        for item in value:
                            st.write(f"- {item}")
                
                # Competitive Analysis
                st.subheader("Competitive Analysis")
                competitive = analysis_result.get("competitive_analysis", {})
                for key, value in competitive.items():
                    if isinstance(value, list) and value:
                        st.write(f"**{key.replace('_', ' ').title()}:**")
                        for item in value:
                            st.write(f"- {item}")
                
                # Risk Analysis
                st.subheader("Risk Analysis")
                risk = analysis_result.get("risk_analysis", {})
                for key, value in risk.items():
                    if isinstance(value, list) and value:
                        st.write(f"**{key.replace('_', ' ').title()}:**")
                        for item in value:
                            st.write(f"- {item}")
                
                # Recommendations
                st.subheader("Recommendations")
                recommendations = analysis_result.get("recommendations", {})
                for key, value in recommendations.items():
                    if isinstance(value, list) and value:
                        st.write(f"**{key.replace('_', ' ').title()}:**")
                        for item in value:
                            st.write(f"- {item}")
                
                # Generate and display visualizations
                st.header("Visualizations")
                
                # Create presentation
                with st.spinner("Generating presentation..."):
                    pptx_path = "patent_analysis.pptx"
                    visualizer.create_presentation(selected_patent, analysis_result, pptx_path)
                    with open(pptx_path, "rb") as f:
                        st.download_button(
                            "Download Presentation",
                            f,
                            file_name="patent_analysis.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                
                # Create flowchart
                with st.spinner("Generating flowchart..."):
                    flowchart_path = "patent_flowchart"
                    visualizer.create_flowchart(selected_patent, flowchart_path)
                    with open(f"{flowchart_path}.pdf", "rb") as f:
                        st.download_button(
                            "Download Flowchart",
                            f,
                            file_name="patent_flowchart.pdf",
                            mime="application/pdf"
                        )
                
                # Create radar chart
                with st.spinner("Generating radar chart..."):
                    radar_path = "patent_radar.png"
                    visualizer.create_radar_chart(analysis_result, radar_path)
                    with open(radar_path, "rb") as f:
                        st.download_button(
                            "Download Radar Chart",
                            f,
                            file_name="patent_radar.png",
                            mime="image/png"
                        )
        
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()
